#!/usr/bin/env python3
"""PPTX Generator — create PowerPoint presentations from JSON or Markdown.

Usage:
  python3 pptx_tool.py generate --input slides.json --output out.pptx
  python3 pptx_tool.py from-markdown --input content.md --output out.pptx
  python3 pptx_tool.py add-slide --file existing.pptx --title "Title" --content "Body"
"""
from __future__ import annotations
import argparse, json, os, re, sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
except ImportError:
    sys.exit("python-pptx not installed. Run: pip3 install python-pptx")


DEFAULT_OUTPUT_DIR = os.path.join(
    os.environ.get("OPENCLAW_HOME", os.path.expanduser("~/.xyvaclaw")),
    "workspace", "output", "pptx"
)


def _ensure_dir(path: str):
    os.makedirs(os.path.dirname(path) or ".", exist_ok=True)


def _apply_font(run, size=18, bold=False, color=None):
    run.font.size = Pt(size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = RGBColor(*color)


def _add_title_slide(prs, title: str, subtitle: str = ""):
    layout = prs.slide_layouts[0]  # Title Slide
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    if subtitle and slide.placeholders[1]:
        slide.placeholders[1].text = subtitle
    return slide


def _add_content_slide(prs, title: str, bullets: list[str]):
    layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for i, bullet in enumerate(bullets):
        if i == 0:
            tf.text = bullet
        else:
            p = tf.add_paragraph()
            p.text = bullet
        # Style
        para = tf.paragraphs[i] if i < len(tf.paragraphs) else tf.paragraphs[-1]
        para.level = 0
        for run in para.runs:
            run.font.size = Pt(18)
    return slide


def _add_table_slide(prs, title: str, headers: list[str], rows: list[list[str]]):
    layout = prs.slide_layouts[5]  # Blank
    slide = prs.slides.add_slide(layout)
    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    _apply_font(p.runs[0] if p.runs else p.add_run(), size=28, bold=True)

    # Table
    num_rows = len(rows) + 1
    num_cols = len(headers)
    left, top, width, height = Inches(0.5), Inches(1.3), Inches(9), Inches(0.4 * num_rows)
    table = slide.shapes.add_table(num_rows, num_cols, left, top, width, height).table

    # Headers
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                run.font.bold = True
                run.font.size = Pt(14)

    # Data rows
    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(12)

    return slide


def _add_image_slide(prs, title: str, image_path: str):
    layout = prs.slide_layouts[5]  # Blank
    slide = prs.slides.add_slide(layout)
    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    _apply_font(p.runs[0] if p.runs else p.add_run(), size=28, bold=True)

    # Image
    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(1), Inches(1.5), Inches(8), Inches(5))
    else:
        txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1))
        txBox2.text_frame.text = f"[Image not found: {image_path}]"

    return slide


def generate_from_json(data: dict, output_path: str):
    """Generate PPTX from a JSON structure."""
    prs = Presentation()

    for slide_def in data.get("slides", []):
        layout = slide_def.get("layout", "content")
        title = slide_def.get("title", "")

        if layout == "title":
            _add_title_slide(prs, title, slide_def.get("subtitle", ""))
        elif layout == "content":
            _add_content_slide(prs, title, slide_def.get("bullets", []))
        elif layout == "table":
            _add_table_slide(prs, title,
                           slide_def.get("headers", []),
                           slide_def.get("rows", []))
        elif layout == "image":
            _add_image_slide(prs, title, slide_def.get("image", ""))
        else:
            _add_content_slide(prs, title, slide_def.get("bullets", [title]))

    _ensure_dir(output_path)
    prs.save(output_path)
    print(f"✅ Generated: {output_path} ({len(prs.slides)} slides)")
    return output_path


def generate_from_markdown(md_path: str, output_path: str):
    """Generate PPTX from Markdown — each ## heading becomes a slide."""
    with open(md_path, "r", encoding="utf-8") as f:
        content = f.read()

    slides = []
    # Split by ## headings
    sections = re.split(r'^## ', content, flags=re.MULTILINE)

    # First section might have a # title
    first = sections[0].strip()
    title_match = re.match(r'^# (.+)', first, re.MULTILINE)
    if title_match:
        slides.append({"layout": "title", "title": title_match.group(1).strip()})

    for section in sections[1:]:
        lines = section.strip().split('\n')
        heading = lines[0].strip()
        bullets = []
        for line in lines[1:]:
            line = line.strip()
            if line.startswith('- ') or line.startswith('* '):
                bullets.append(line[2:])
            elif line and not line.startswith('#'):
                bullets.append(line)
        if bullets:
            slides.append({"layout": "content", "title": heading, "bullets": bullets})
        else:
            slides.append({"layout": "content", "title": heading, "bullets": [heading]})

    return generate_from_json({"slides": slides}, output_path)


def add_slide_to_existing(pptx_path: str, title: str, content: str):
    """Add a slide to an existing PPTX file."""
    prs = Presentation(pptx_path)
    bullets = [line.strip() for line in content.split('\n') if line.strip()]
    _add_content_slide(prs, title, bullets)
    prs.save(pptx_path)
    print(f"✅ Added slide to: {pptx_path} (now {len(prs.slides)} slides)")


def main():
    parser = argparse.ArgumentParser(description="PPTX Generator")
    sub = parser.add_subparsers(dest="cmd")

    # generate
    gen = sub.add_parser("generate", help="Generate PPTX from JSON")
    gen.add_argument("--input", "-i", required=True, help="JSON file path")
    gen.add_argument("--output", "-o", default=None, help="Output PPTX path")

    # from-markdown
    md = sub.add_parser("from-markdown", help="Generate PPTX from Markdown")
    md.add_argument("--input", "-i", required=True, help="Markdown file path")
    md.add_argument("--output", "-o", default=None, help="Output PPTX path")

    # add-slide
    add = sub.add_parser("add-slide", help="Add slide to existing PPTX")
    add.add_argument("--file", "-f", required=True, help="Existing PPTX path")
    add.add_argument("--title", "-t", required=True, help="Slide title")
    add.add_argument("--content", "-c", required=True, help="Slide content")

    args = parser.parse_args()

    if not args.cmd:
        parser.print_help()
        return

    if args.cmd == "generate":
        with open(args.input, "r", encoding="utf-8") as f:
            data = json.load(f)
        output = args.output or os.path.join(DEFAULT_OUTPUT_DIR, "output.pptx")
        generate_from_json(data, output)

    elif args.cmd == "from-markdown":
        output = args.output or os.path.join(DEFAULT_OUTPUT_DIR, "output.pptx")
        generate_from_markdown(args.input, output)

    elif args.cmd == "add-slide":
        add_slide_to_existing(args.file, args.title, args.content)


if __name__ == "__main__":
    main()
